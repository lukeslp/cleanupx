"""
Content extraction module for Reference Renamer.
Handles extraction of text and metadata from various file formats.
"""

from pathlib import Path
from typing import Dict, Any, Optional
import logging
import io
import tempfile
import re
import os
import shutil

import PyPDF2
import pdfplumber
import docx
import pptx
import markdown
import textract
from tika import parser as tika_parser
from pdfminer.high_level import extract_text as pdfminer_extract
from PIL import Image
import pytesseract
from pdf2image import convert_from_path
from pdf2image.exceptions import PDFPageCountError

from ..utils.exceptions import ContentExtractionError, ExtractionError
from ..utils.logging import get_logger

class ContentExtractor:
    """Extracts text content from various file formats."""
    
    def __init__(self, logger: Optional[logging.Logger] = None):
        """
        Initialize the content extractor.
        
        Args:
            logger: Optional logger instance
        """
        self.logger = logger or get_logger(__name__)
        
    def extract_text(self, file_path: Path) -> Dict[str, str]:
        """Extract text from a document using the appropriate method.
        
        Args:
            file_path: Path to the document
            
        Returns:
            Dictionary containing extracted text sections
            
        Raises:
            ExtractionError: If text extraction fails
        """
        try:
            suffix = file_path.suffix.lower()
            
            # Try tika first for any format
            try:
                tika_content = self._extract_with_tika(file_path)
                if tika_content and len(tika_content.strip()) > 100:  # Reasonable content found
                    return {'full_text': tika_content}
            except Exception as e:
                self.logger.warning(f"Tika extraction failed, falling back to specific extractors: {e}")
            
            # Fall back to specific extractors
            if suffix == '.pdf':
                return self._extract_pdf_text(file_path)
            elif suffix == '.docx':
                return self._extract_docx_text(file_path)
            elif suffix == '.pptx':
                return self._extract_pptx_text(file_path)
            elif suffix == '.md':
                return self._extract_markdown_text(file_path)
            elif suffix == '.txt':
                return self._extract_text_file(file_path)
            else:
                # Try textract as a last resort
                try:
                    content = textract.process(str(file_path)).decode('utf-8')
                    return {'full_text': content}
                except Exception as e:
                    raise ExtractionError(f"Unsupported file format: {suffix}")
                    
        except Exception as e:
            raise ExtractionError(f"Failed to extract text from {file_path}: {str(e)}")
            
    def _extract_with_tika(self, file_path: Path) -> str:
        """Extract text using Apache Tika."""
        parsed = tika_parser.from_file(str(file_path))
        return parsed.get('content', '')
            
    def _extract_pdf_text(self, file_path: Path) -> Dict[str, str]:
        """Extract text from PDF using multiple methods."""
        content = {'full_text': '', 'header': '', 'first_page': ''}
        
        # Try pdfplumber first
        try:
            with pdfplumber.open(file_path) as pdf:
                text = ''
                for i, page in enumerate(pdf.pages):
                    page_text = page.extract_text()
                    if page_text:
                        if i == 0:
                            content['first_page'] = page_text
                            # Extract header (first few lines)
                            lines = page_text.split('\\n')[:3]
                            content['header'] = '\\n'.join(lines)
                        text += page_text + '\\n'
                content['full_text'] = text
                
                if text.strip():  # If we got good content, return it
                    return content
        except Exception as e:
            self.logger.warning(f"pdfplumber extraction failed: {e}")
        
        # Try pdfminer if pdfplumber failed
        try:
            text = pdfminer_extract(str(file_path))
            if text.strip():
                content['full_text'] = text
                lines = text.split('\\n')
                content['header'] = '\\n'.join(lines[:3])
                content['first_page'] = '\\n'.join(lines[:20])  # Approximate first page
                return content
        except Exception as e:
            self.logger.warning(f"pdfminer extraction failed: {e}")
        
        # Try PyPDF2 as last resort
        try:
            with open(file_path, 'rb') as file:
                reader = PyPDF2.PdfReader(file)
                text = ""
                for i, page in enumerate(reader.pages):
                    page_text = page.extract_text()
                    if not page_text.strip():
                        # Try OCR if no text found
                        page_text = self._ocr_pdf_page(file_path, i)
                    if i == 0:
                        content['first_page'] = page_text
                        lines = page_text.split('\\n')[:3]
                        content['header'] = '\\n'.join(lines)
                    text += page_text + "\\n"
                content['full_text'] = text
                return content
        except Exception as e:
            self.logger.error(f"All PDF extraction methods failed: {e}")
            raise ExtractionError(f"Could not extract text from PDF: {e}")
            
    def _extract_docx_text(self, file_path: Path) -> Dict[str, str]:
        """Extract text from DOCX file."""
        try:
            doc = docx.Document(file_path)
            content = {'full_text': '', 'header': '', 'first_page': ''}
            
            # Extract all paragraphs
            text = '\\n'.join(p.text for p in doc.paragraphs)
            content['full_text'] = text
            
            # Get header and first page
            first_paras = [p.text for p in doc.paragraphs[:20]]  # Approximate first page
            content['header'] = '\\n'.join(first_paras[:3])
            content['first_page'] = '\\n'.join(first_paras)
            
            return content
        except Exception as e:
            raise ExtractionError(f"Failed to extract text from DOCX: {e}")
            
    def _extract_pptx_text(self, file_path: Path) -> Dict[str, str]:
        """Extract text from PPTX file."""
        try:
            prs = pptx.Presentation(file_path)
            content = {'full_text': '', 'header': '', 'first_page': ''}
            
            all_text = []
            first_slide_text = []
            
            for i, slide in enumerate(prs.slides):
                slide_text = []
                for shape in slide.shapes:
                    if hasattr(shape, "text"):
                        slide_text.append(shape.text)
                
                if i == 0:
                    first_slide_text = slide_text
                all_text.extend(slide_text)
            
            content['full_text'] = '\\n'.join(all_text)
            content['header'] = '\\n'.join(first_slide_text[:3])
            content['first_page'] = '\\n'.join(first_slide_text)
            
            return content
        except Exception as e:
            raise ExtractionError(f"Failed to extract text from PPTX: {e}")
            
    def _extract_markdown_text(self, file_path: Path) -> Dict[str, str]:
        """Extract text from Markdown file."""
        try:
            with open(file_path, 'r', encoding='utf-8') as f:
                md_text = f.read()
                
            # Convert to HTML and extract text
            html = markdown.markdown(md_text)
            # Simple HTML tag removal (could use BeautifulSoup for better parsing)
            text = re.sub(r'<[^>]+>', '', html)
            
            content = {'full_text': text}
            
            # Extract header and first page
            lines = text.split('\\n')
            content['header'] = '\\n'.join(lines[:3])
            content['first_page'] = '\\n'.join(lines[:20])  # Approximate first page
            
            return content
        except Exception as e:
            raise ExtractionError(f"Failed to extract text from Markdown: {e}")
            
    def _extract_text_file(self, file_path: Path) -> Dict[str, str]:
        """Extract text from plain text file."""
        try:
            with open(file_path, 'r', encoding='utf-8') as f:
                text = f.read()
                
            content = {'full_text': text}
            
            # Extract header and first page
            lines = text.split('\\n')
            content['header'] = '\\n'.join(lines[:3])
            content['first_page'] = '\\n'.join(lines[:20])  # Approximate first page
            
            return content
        except UnicodeDecodeError:
            # Try different encodings
            for encoding in ['latin-1', 'cp1252', 'iso-8859-1']:
                try:
                    with open(file_path, 'r', encoding=encoding) as f:
                        text = f.read()
                    content = {
                        'full_text': text,
                        'header': '\\n'.join(text.split('\\n')[:3]),
                        'first_page': '\\n'.join(text.split('\\n')[:20])
                    }
                    return content
                except UnicodeDecodeError:
                    continue
            raise ExtractionError("Failed to decode text file with any encoding")
            
    def _ocr_pdf_page(self, file_path: Path, page_num: int) -> str:
        """Perform OCR on a PDF page."""
        temp_dir = None
        try:
            temp_dir = tempfile.mkdtemp()
            
            # Convert PDF page to image
            images = convert_from_path(file_path, first_page=page_num+1, last_page=page_num+1)
            if not images:
                return ""
                
            # Perform OCR on the image
            text = pytesseract.image_to_string(images[0])
            return text.strip()
            
        except Exception as e:
            self.logger.error(f"OCR failed: {e}")
            return ""
            
        finally:
            if temp_dir and os.path.exists(temp_dir):
                shutil.rmtree(temp_dir)

    def extract_images(self, file_path: Path) -> Dict[str, Any]:
        """
        Extracts images from a document if present.
        
        Args:
            file_path: Path to the document
            
        Returns:
            Dict containing:
                - images: List of extracted images
                - count: Number of images found
        """
        if file_path.suffix.lower() != '.pdf':
            return {'images': [], 'count': 0}
            
        try:
            with open(file_path, 'rb') as file:
                reader = PyPDF2.PdfReader(file)
                images = []
                
                for i, page in enumerate(reader.pages):
                    if '/XObject' in page['/Resources']:
                        xObject = page['/Resources']['/XObject'].get_object()
                        
                        for obj in xObject:
                            if xObject[obj]['/Subtype'] == '/Image':
                                self.logger.debug(f"Found image on page {i+1}")
                                images.append({
                                    'page': i + 1,
                                    'size': (xObject[obj]['/Width'], xObject[obj]['/Height']),
                                    'format': xObject[obj]['/Filter']
                                })
                                
                return {
                    'images': images,
                    'count': len(images)
                }
                
        except Exception as e:
            self.logger.error(f"Error extracting images: {str(e)}")
            return {'images': [], 'count': 0}

    def _clean_text(self, text: str) -> str:
        """Clean and optimize text for processing.
        
        Args:
            text: Raw text to clean
            
        Returns:
            Cleaned text
        """
        # Remove excessive whitespace
        text = ' '.join(text.split())
        
        # Remove common PDF artifacts
        text = re.sub(r'\f', ' ', text)  # Form feed
        text = re.sub(r'(?<!\n)\n(?!\n)', ' ', text)  # Single newlines
        text = re.sub(r'\n{3,}', '\n\n', text)  # Multiple newlines
        
        # Remove header/footer artifacts
        text = re.sub(r'(?m)^\d+$', '', text)  # Page numbers
        text = re.sub(r'(?m)^.*?Page \d+.*?$', '', text)  # Page headers
        
        # Clean up special characters
        text = re.sub(r'[^\x00-\x7F]+', ' ', text)  # Remove non-ASCII
        text = re.sub(r'[\x00-\x08\x0B\x0C\x0E-\x1F\x7F]', '', text)  # Remove control chars
        
        # Remove repeated punctuation
        text = re.sub(r'([.,!?])\1+', r'\1', text)
        
        # Normalize spaces
        text = re.sub(r'\s+', ' ', text)
        
        return text.strip()

    def extract_text(self, file_path: Path) -> Dict[str, str]:
        """Extract text from a document.
        
        Args:
            file_path: Path to the document
            
        Returns:
            Dictionary containing extracted text sections
            
        Raises:
            ExtractionError: If text extraction fails
        """
        try:
            self.logger.info(f"Extracting text from {file_path}")
            
            # Read PDF
            with open(file_path, 'rb') as file:
                reader = PyPDF2.PdfReader(file)
                self.logger.info(f"PDF has {len(reader.pages)} pages")
                
                if len(reader.pages) == 0:
                    raise ExtractionError("PDF has no pages")
                
                # Initialize content sections
                content = {
                    'header': '',
                    'title_page': '',
                    'first_page': '',
                    'references': '',
                    'full_text': ''
                }
                
                # Extract first page (header and title)
                first_page_text = ""
                try:
                    first_page_text = reader.pages[0].extract_text()
                    self.logger.info(f"Raw first page text: {first_page_text[:500] if first_page_text else ''}")
                except Exception as e:
                    self.logger.error(f"Failed to extract text from first page: {str(e)}")
                
                # Attempt OCR if no text found
                if not first_page_text.strip():
                    self.logger.warning("No text found in PDF, attempting OCR")
                    try:
                        first_page_text = self._ocr_pdf_page(file_path, 0)
                        self.logger.info(f"OCR result: {first_page_text[:500] if first_page_text else ''}")
                    except Exception as e:
                        self.logger.error(f"OCR failed for first page: {str(e)}")
                
                # Clean and process first page text
                if first_page_text:
                    first_page_text = self._clean_text(first_page_text)
                    self.logger.info(f"Cleaned first page text: {first_page_text[:500]}")
                    
                    # Split first page into sections
                    lines = first_page_text.split('\n')
                    content['header'] = '\n'.join(lines[:3])  # First 3 lines usually contain title
                    content['title_page'] = first_page_text
                    content['first_page'] = first_page_text
                
                # Extract references from last few pages
                references_text = []
                last_pages = min(3, len(reader.pages))
                for i in range(last_pages):
                    page_num = len(reader.pages) - 1 - i
                    try:
                        page_text = reader.pages[page_num].extract_text()
                        if page_text:
                            self.logger.info(f"Page {page_num} text: {page_text[:200]}")
                            
                            # Look for references section
                            if any(marker in page_text.lower() for marker in ['references', 'bibliography', 'works cited']):
                                page_text = self._clean_text(page_text)
                                references_text.append(page_text)
                                self.logger.info(f"Found references on page {page_num}")
                    except Exception as e:
                        self.logger.error(f"Failed to extract text from page {page_num}: {str(e)}")
                        continue
                
                if references_text:
                    content['references'] = '\n'.join(references_text)
                
                # Extract full text (first 5 pages)
                full_text = []
                max_pages = min(5, len(reader.pages))
                for i in range(max_pages):
                    try:
                        page_text = reader.pages[i].extract_text()
                        if not page_text.strip():
                            self.logger.info(f"No text on page {i}, attempting OCR")
                            page_text = self._ocr_pdf_page(file_path, i)
                        
                        if page_text:
                            page_text = self._clean_text(page_text)
                            full_text.append(page_text)
                            self.logger.info(f"Page {i} length: {len(page_text)}")
                    except Exception as e:
                        self.logger.error(f"Failed to extract text from page {i}: {str(e)}")
                        continue
                
                if full_text:
                    content['full_text'] = '\n'.join(full_text)
                
                # Log content summary
                for section, text in content.items():
                    self.logger.info(f"{section} length: {len(text)}")
                    if text:
                        self.logger.info(f"{section} preview: {text[:200]}")
                
                return content
                
        except Exception as e:
            error_msg = f"Failed to extract text from {file_path}: {str(e)}"
            self.logger.error(error_msg)
            raise ExtractionError(error_msg)
    
    def _find_references_section(self, pages: list[str]) -> Optional[str]:
        """Find and extract the references section from the document.
        
        Args:
            pages: List of page texts
            
        Returns:
            References section text if found, None otherwise
        """
        references_text = []
        in_references = False
        
        for page in pages:
            # Look for references section header
            if any(header in page.lower() for header in ['references', 'bibliography', 'works cited']):
                in_references = True
                # Extract text after the header
                start_idx = max(page.lower().find('references'),
                              page.lower().find('bibliography'),
                              page.lower().find('works cited'))
                references_text.append(page[start_idx:])
            elif in_references:
                # Continue collecting text until end of document or new section
                if any(header in page.lower() for header in ['appendix', 'acknowledgments', 'notes']):
                    break
                references_text.append(page)
        
        return '\n'.join(references_text) if references_text else None 